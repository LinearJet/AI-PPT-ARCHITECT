# presentation_exporter.py

import asyncio
import io
from playwright.async_api import async_playwright
from pptx import Presentation
from pptx.util import Inches
from pptx.dml.color import RGBColor
from bs4 import BeautifulSoup

# --- Constants ---
VIEWPORT_WIDTH_PX = 1280
VIEWPORT_HEIGHT_PX = 720
PPTX_WIDTH_INCHES = 13.333  # 1280px / 96 DPI
PPTX_HEIGHT_INCHES = 7.5     # 720px / 96 DPI

class StaticImageExporter:
    """
    Exports a presentation by taking high-resolution screenshots of each slide
    and compiling them into a PPTX file.
    """
    def __init__(self, slides_html):
        self.slides_html = slides_html

    def _hex_to_rgb(self, hex_color):
        """Converts a hex color string to an (R, G, B) tuple."""
        hex_color = hex_color.lstrip('#')
        if len(hex_color) == 3:
            hex_color = hex_color * 2
        return tuple(int(hex_color[i:i+2], 16) for i in (0, 2, 4))

    async def export(self):
        """
        Main public method to run the entire export process.
        - Launches a headless browser.
        - Takes a screenshot of each slide's HTML.
        - Compiles the screenshots into a PPTX file.
        """
        print("Exporter: Initializing headless browser for static export...")
        screenshots = []

        async with async_playwright() as p:
            browser = await p.chromium.launch()
            page = await browser.new_page()
            await page.set_viewport_size({"width": VIEWPORT_WIDTH_PX, "height": VIEWPORT_HEIGHT_PX})

            num_slides = len(self.slides_html)
            for i, html_content in enumerate(self.slides_html):
                print(f" - Capturing screenshot for Slide {i + 1}/{num_slides}...")
                await page.set_content(html_content)
                # --- FIX 2: Extend timeout for animations to complete ---
                await page.wait_for_timeout(5000) 
                
                screenshot_bytes = await page.screenshot(type='png')
                screenshots.append(io.BytesIO(screenshot_bytes))

            await browser.close()

        print("Exporter: All slides captured. Compiling PPTX file...")
        prs = Presentation()
        prs.slide_width = Inches(PPTX_WIDTH_INCHES)
        prs.slide_height = Inches(PPTX_HEIGHT_INCHES)

        for screenshot_buffer in screenshots:
            slide_layout = prs.slide_layouts[6]  # Blank layout
            slide = prs.slides.add_slide(slide_layout)
            
            slide.shapes.add_picture(
                screenshot_buffer,
                left=Inches(0),
                top=Inches(0),
                width=prs.slide_width
            )

        output_buffer = io.BytesIO()
        prs.save(output_buffer)
        output_buffer.seek(0)
        
        print("Exporter: Static PPTX compilation complete!")
        return output_buffer.getvalue()